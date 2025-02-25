import streamlit as st
from pptx import Presentation
import tempfile

st.set_page_config(page_title="🔄 Générateur de Présentation PowerPoint", layout="centered")

st.title("🔄 Générateur de Présentation PowerPoint")
st.write("📤 **Importe un fichier PowerPoint et remplis les champs pour le personnaliser.**")

# 📌 Upload du fichier PowerPoint
uploaded_file = st.file_uploader("📤 **Importer un fichier .pptx**", type=["pptx"])

if uploaded_file:
    with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp_file:
        tmp_file.write(uploaded_file.read())
        tmp_path = tmp_file.name

    st.success("✅ Fichier PowerPoint importé avec succès !")

    # 📌 Champs de saisie des valeurs
    st.subheader("✍️ Remplis les informations suivantes :")

    nom_entreprise = st.text_input("🏢 Nom de l'entreprise")
    secteur_entreprise = st.text_input("💼 Secteur d'activité")
    responsable_entreprise = st.text_input("👤 Responsable")
    localisation_entreprise = st.text_input("📍 Localisation")
    effectif_entreprise = st.text_input("👥 Effectif")
    date_entreprise = st.text_input("📅 Date d’exercice")
    responsable_cabinet = st.text_input("📑 Responsable du Cabinet")

    prix_comptabilite = st.text_input("💰 Prix Comptabilité (€)")
    prix_fiscalite = st.text_input("💰 Prix Fiscalité (€)")
    prix_audit = st.text_input("💰 Prix Audit (€)")
    prix_total = st.text_input("💰 Total (€)")
    prix_devis = st.text_input("💰 Montant du Devis (€)")
    prix_devis_mois = st.text_input("💰 Montant mensuel TTC (€)")
    nombre_salaries = st.text_input("👥 Nombre de Salariés")
    prix_dossier = st.text_input("💰 Prix Configuration du Dossier (€)")
    prix_bulletin = st.text_input("💰 Prix d’un Bulletin (€)")
    prix_general = st.text_input("💰 Prix Général (€)")

    replacements = {
        "NomEntreprise": nom_entreprise,
        "SecteurEntreprise": secteur_entreprise,
        "ResponsableEntreprise": responsable_entreprise,
        "LocalisationEntreprise": localisation_entreprise,
        "EffectifEntreprise": effectif_entreprise,
        "DateEntreprise": date_entreprise,
        "ResponsableCabinet": responsable_cabinet,
        "PrixComptabilité": prix_comptabilite,
        "PrixFiscalité": prix_fiscalite,
        "PrixAudit": prix_audit,
        "PrixTotal": prix_total,
        "PrixDevis": prix_devis,
        "PrisDevisMois": prix_devis_mois,
        "Nombresalariés": nombre_salaries,
        "PrixDossier": prix_dossier,
        "PrixBulletin": prix_bulletin,
        "PrixFinale": prix_general,
    }

    def replace_text_in_shapes(shapes, replacements):
        """Remplace le texte dans les formes et tableaux."""
        for shape in shapes:
            try:
                if shape.has_text_frame and shape.text_frame is not None:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            for key, value in replacements.items():
                                if key in run.text and value:
                                    run.text = run.text.replace(key, value)

                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            for key, value in replacements.items():
                                if key in cell.text and value:
                                    cell.text = cell.text.replace(key, value)

                if hasattr(shape, "shapes"):
                    replace_text_in_shapes(shape.shapes, replacements)

            except Exception as e:
                st.error(f"⚠️ Erreur sur une forme : {e}")

    def update_pptx(input_pptx, replacements):
        prs = Presentation(input_pptx)
        for slide in prs.slides:
            replace_text_in_shapes(slide.shapes, replacements)
        
        output_pptx = "presentation_modifiee.pptx"
        prs.save(output_pptx)
        return output_pptx

    if st.button("🛠️ Générer le PowerPoint"):
        output_file = update_pptx(tmp_path, replacements)
        with open(output_file, "rb") as f:
            st.download_button("📥 Télécharger la présentation modifiée", f, file_name="presentation_modifiee.pptx")
        st.success("✅ Présentation mise à jour et prête au téléchargement !")
